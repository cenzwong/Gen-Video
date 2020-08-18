#!/usr/bin/env python
# coding: utf-8

from pptx import Presentation
import os

from bs4 import BeautifulSoup
import requests
from gtts import gTTS
from mutagen.mp3 import MP3

import pickle


myInput = input("https://ctext.org/...")
print("you enter:" + myInput)
url = "https://ctext.org/" + myInput


# In[3]:


prs = Presentation()
layout = prs.slide_layouts[0]
# url = "https://ctext.org/mengzi/jin-xin-ii/zh"

html = requests.get(url)
bsObj = BeautifulSoup(html.text,"lxml")


# In[18]:


name = bsObj.find("meta", {"property":"og:title"}).get("content")
filename = name + ".pptx"


# In[46]:


temp = bsObj.find("div", {"id":"content3"})
myTable = temp.findAll("table")[1]
myList = myTable.findAll("td", {"class","ctext"})


language = "zh-tw"
FolderName = os.path.join("Python_audio_v2", "output")

duration_list = []
myContent = []
for i, t in enumerate(myList):
    mysentence = t.get_text()
    mysentence = mysentence.replace('\n', '')
    if mysentence == str(''):
        continue
    myContent.append(mysentence)

## Save Powerpoint file
def create_slides_from_list(content):
    for s in content:
        slide = prs.slides.add_slide(layout)
        slide.placeholders[1].text = s

create_slides_from_list(myContent)

prs.save(filename)
# os.startfile(filename)

## Save MP3
for i, mysentence in enumerate(myContent):
    # valid sentence

    tts = gTTS(text=mysentence, lang=language, slow=False)
    
    mp3_filename = str(i).zfill(4) + ".mp3"
    tts.save(os.path.join(FolderName, mp3_filename))
    
    audio = MP3(os.path.join(FolderName, mp3_filename))
    duration_list.append(audio.info.length)
    print("Processing: "+ str(round(i*100/len(myList), 2)) +"% finish", end='\r')
# myContent


# In[47]:
with open(os.path.join("Python_video_v2", "mylist"), 'wb') as f: 
    pickle.dump(duration_list, f)

with open(os.path.join("Video_combine", "video_name"), 'wb') as f: 
    pickle.dump(name, f)
# with open('mylist', 'rb') as f: 
#     mylist = pickle.load(f) 



